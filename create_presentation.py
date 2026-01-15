#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
KooRemapper 발표자료 생성 프로그램
README.md를 파싱하여 PowerPoint 프레젠테이션을 생성합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        # 색상 테마
        self.primary_color = RGBColor(31, 78, 120)  # 다크 블루
        self.accent_color = RGBColor(68, 114, 196)  # 블루
        self.text_color = RGBColor(50, 50, 50)      # 다크 그레이
        
    def add_title_slide(self, title, subtitle):
        """타이틀 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # 배경 (그라디언트 효과 대신 단색)
        background = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.primary_color
        background.line.fill.background()
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 서브타이틀
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(14), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def add_content_slide(self, title, content_items, layout_type='bullet'):
        """콘텐츠 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # 타이틀 배경
        title_bg = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.prs.slide_width,
            Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.primary_color
        title_bg.line.fill.background()
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 콘텐츠
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(14.4), Inches(6.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for idx, item in enumerate(content_items):
            if idx > 0:
                text_frame.add_paragraph()
            
            p = text_frame.paragraphs[idx]
            
            # 레벨 감지
            level = 0
            cleaned_item = item.lstrip('- ')
            if item.startswith('  - '):
                level = 1
                cleaned_item = item[4:]
            elif item.startswith('    - '):
                level = 2
                cleaned_item = item[6:]
            elif item.startswith('- '):
                level = 0
                cleaned_item = item[2:]
            else:
                cleaned_item = item
            
            p.text = cleaned_item
            p.level = level
            p.font.size = Pt(20) if level == 0 else Pt(18)
            p.font.color.rgb = self.text_color
            p.space_after = Pt(10)
            
    def add_code_slide(self, title, code_lines, language='bash'):
        """코드 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 타이틀 배경
        title_bg = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.primary_color
        title_bg.line.fill.background()
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # 코드 박스 배경
        code_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.8),
            Inches(14.4), Inches(6)
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(40, 44, 52)  # 다크 배경
        code_bg.line.color.rgb = RGBColor(100, 100, 100)
        
        # 코드 텍스트
        code_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(14), Inches(5.5)
        )
        text_frame = code_box.text_frame
        text_frame.word_wrap = True
        
        code_text = '\n'.join(code_lines)
        text_frame.text = code_text
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = 'Consolas'
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(171, 178, 191)  # 밝은 그레이
            
    def add_diagram_slide(self, title, diagram_lines):
        """다이어그램 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 타이틀 배경
        title_bg = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.primary_color
        title_bg.line.fill.background()
        
        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # 다이어그램 박스
        diagram_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8),
            Inches(13), Inches(6)
        )
        text_frame = diagram_box.text_frame
        text_frame.word_wrap = True
        
        diagram_text = '\n'.join(diagram_lines)
        text_frame.text = diagram_text
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = 'Consolas'
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_color
            
    def parse_readme(self, readme_path):
        """README.md 파싱 및 슬라이드 생성"""
        with open(readme_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        lines = content.split('\n')
        
        # 타이틀 슬라이드
        self.add_title_slide(
            "KooRemapper",
            "LS-DYNA K-file용 메쉬 매핑 및 응력 분석 도구"
        )
        
        # 주요 기능 추출
        features = []
        in_features = False
        for line in lines:
            if line.strip() == '## 주요 기능':
                in_features = True
                continue
            if in_features and line.startswith('##'):
                break
            if in_features and line.strip().startswith('- '):
                features.append(line.strip())
        
        if features:
            self.add_content_slide("주요 기능", features)
        
        # 워크플로우 다이어그램
        workflow = []
        in_workflow = False
        in_code_block = False
        for line in lines:
            if '## 전체 워크플로우' in line:
                in_workflow = True
                continue
            if in_workflow and line.strip() == '```':
                if not in_code_block:
                    in_code_block = True
                    continue
                else:
                    break
            if in_workflow and in_code_block:
                workflow.append(line)
        
        if workflow:
            self.add_diagram_slide("전체 워크플로우", workflow)
        
        # 명령어별 슬라이드
        self.create_command_slides(lines)
        
        # 예제 슬라이드
        self.create_example_slides(lines)
        
        # 기술 세부사항
        self.create_technical_slides(lines)
        
        # 마지막 슬라이드
        self.add_title_slide(
            "감사합니다",
            "KooRemapper - 메쉬 매핑 및 응력 분석 도구"
        )
    
    def create_command_slides(self, lines):
        """명령어 설명 슬라이드 생성"""
        
        # generate-var 명령어
        self.add_content_slide(
            "1. 가변 메쉬 생성 (generate-var)",
            [
                "YAML 설정으로 곡선/가변 밀도 메쉬 생성",
                "",
                "타입 1: 가변 밀도 평면 메쉬 (type: flat)",
                "  - 영역별로 요소 밀도 제어",
                "  - 밴딩 영역은 조밀, 평평한 영역은 성기게",
                "  - 레퍼런스 기반 자동 크기 조정",
                "",
                "타입 2: 곡선 메쉬 (type: curved)",
                "  - 2D centerline 기반 벤딩 메쉬 생성",
                "  - B-spline, Catmull-Rom 보간",
                "  - 자동 arc-length 계산"
            ]
        )
        
        # generate-var 예제 코드
        yaml_example = [
            "# 곡선 메쉬 YAML 예제",
            "type: curved",
            "centerline_points:",
            "  - [0, 0]",
            "  - [40, 10]",
            "  - [80, 30]",
            "interpolation: catmull_rom",
            "",
            "# 명령어",
            "KooRemapper generate-var config.yaml output.k"
        ]
        self.add_code_slide("generate-var 예제", yaml_example, 'yaml')
        
        # map 명령어
        self.add_content_slide(
            "2. 메쉬 매핑 (map)",
            [
                "디테일 플랫 → 디테일 벤트 변형",
                "",
                "입력:",
                "  - 심플 벤트: 레퍼런스 (형상 정의, ~10K 요소)",
                "  - 디테일 플랫: 매핑 대상 (~1M 요소)",
                "",
                "출력:",
                "  - 디테일 벤트: 변형 결과 (~1M 요소)",
                "",
                "특징:",
                "  - 레퍼런스는 HEX8 정형 메쉬",
                "  - 디테일은 HEX8/TET4 비정형 가능",
                "  - 자동 크기 조정"
            ]
        )
        
        # map 예제
        map_example = [
            "# 심플 레퍼런스 생성 (10,000 요소)",
            "KooRemapper generate-var curved.yaml simple_bent.k",
            "",
            "# 디테일 플랫 준비 (CAD 또는 generate-var)",
            "# detail_flat.k (1,000,000 요소)",
            "",
            "# 매핑: 디테일 플랫 → 디테일 벤트",
            "KooRemapper map simple_bent.k detail_flat.k detail_bent.k"
        ]
        self.add_code_slide("map 예제", map_example)
        
        # prestress 명령어
        self.add_content_slide(
            "3. 초기 응력 계산 (prestress)",
            [
                "변형 전/후 메쉬로부터 응력 계산",
                "",
                "물성 정의 방법:",
                "  - K-file 물성 자동 읽기 (*PART, *MAT_ELASTIC)",
                "  - 명령줄 옵션으로 덮어쓰기 (-E, -nu)",
                "  - 다중 파트 지원",
                "",
                "출력 포맷:",
                "  - dynain: *INITIAL_STRESS_SOLID",
                "  - CSV: 요소별 응력 데이터",
                "",
                "LS-DYNA에서 *INCLUDE로 바로 사용"
            ]
        )
        
        # prestress 예제
        prestress_example = [
            "# K-file 물성 자동 사용",
            "KooRemapper prestress detail_flat.k detail_bent.k output.dynain",
            "",
            "# 명령줄 물성 지정 (덮어쓰기)",
            "KooRemapper prestress -E 210000 -nu 0.3 \\",
            "    detail_flat.k detail_bent.k output.dynain",
            "",
            "# LS-DYNA 입력 파일",
            "*INCLUDE",
            "output.dynain"
        ]
        self.add_code_slide("prestress 예제", prestress_example)
        
        # 기타 명령어
        self.add_content_slide(
            "4. 기타 명령어",
            [
                "unfold - 벤트 → 플랫 언폴딩",
                "  - Arc-length 기반 평면화",
                "  - 크기 참조용 (선택적)",
                "",
                "generate - 테스트 예제 생성",
                "  - arc, torus, helix, wave 등",
                "  - bent/flat 쌍 자동 생성",
                "",
                "info - 메쉬 정보 출력",
                "  - 노드/요소 개수",
                "  - Bounding Box",
                "  - 정형 그리드 구조 (i,j,k)"
            ]
        )
    
    def create_example_slides(self, lines):
        """예제 슬라이드 생성"""
        
        # 전체 예제
        full_example = [
            "# 1. 심플 벤트 레퍼런스 (10K 요소)",
            "KooRemapper generate-var curved.yaml simple_bent.k",
            "",
            "# 2. 심플 플랫 (크기 참조용, 선택)",
            "KooRemapper unfold simple_bent.k simple_flat.k",
            "",
            "# 3. 디테일 플랫 (1M 요소)",
            "KooRemapper generate-var vardens.yaml detail_flat.k",
            "",
            "# 4. 디테일 플랫 → 디테일 벤트",
            "KooRemapper map simple_bent.k detail_flat.k detail_bent.k",
            "",
            "# 5. 초기 응력 계산",
            "KooRemapper prestress detail_flat.k detail_bent.k prestress.dynain"
        ]
        self.add_code_slide("전체 워크플로우 예제", full_example)
    
    def create_technical_slides(self, lines):
        """기술 세부사항 슬라이드 생성"""
        
        self.add_content_slide(
            "기술 세부사항",
            [
                "지원 요소:",
                "  - 레퍼런스: HEX8 (정형 필수)",
                "  - 매핑 대상: HEX8, TET4 (비정형 가능)",
                "",
                "매핑 알고리즘:",
                "  - 정형 그리드 분석 (i,j,k 인덱스)",
                "  - Edge arc-length 계산",
                "  - 파라메트릭 좌표 변환 (u,v,w)",
                "  - Transfinite 보간",
                "",
                "응력 계산:",
                "  - 변형 구배 텐서 (F)",
                "  - Green-Lagrange / Engineering strain",
                "  - Hooke's Law (선형 탄성)",
                "  - von Mises, Principal stress"
            ]
        )
        
        self.add_content_slide(
            "파일 포맷 및 특징",
            [
                "입출력 포맷:",
                "  - K-file: LS-DYNA 메쉬 (*.k)",
                "  - YAML: 설정 파일 (*.yaml)",
                "  - dynain: 초기 응력 (*.dynain)",
                "  - CSV: 스트레인 데이터 (*.csv)",
                "",
                "주요 특징:",
                "  - 크로스 플랫폼 (Windows/Linux/macOS)",
                "  - 정적 링크 (단일 실행파일)",
                "  - K-file 물성 자동 파싱",
                "  - 다중 파트 지원",
                "  - 대용량 메쉬 처리 (100만+ 요소)"
            ]
        )
    
    def save(self, output_path):
        """프레젠테이션 저장"""
        self.prs.save(output_path)
        print(f"✓ 프레젠테이션 생성 완료: {output_path}")

def main():
    print("=" * 60)
    print("KooRemapper 발표자료 생성")
    print("=" * 60)
    
    generator = PresentationGenerator()
    
    print("\n[1/2] README.md 파싱 중...")
    generator.parse_readme("README.md")
    
    print("[2/2] 프레젠테이션 저장 중...")
    generator.save("KooRemapper_Presentation.pptx")
    
    print("\n" + "=" * 60)
    print("완료!")
    print("파일: KooRemapper_Presentation.pptx")
    print("=" * 60)

if __name__ == "__main__":
    main()
